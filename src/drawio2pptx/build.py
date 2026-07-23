"""Turn one draw.io page into individually editable PowerPoint objects."""
from __future__ import annotations

import base64
import io
import re
import shutil
import tempfile
from dataclasses import dataclass, field
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from . import model, stencils
from .svgmap import Label, SvgMap, normalize_color, strip_markup

EMU_PER_PT = 12700
SLIDE_SIZES = {                       # inches -> EMU
    "16:9": (12192000, 6858000),
    "4:3": (9144000, 6858000),
    "16:10": (12192000, 7620000),
}
FRAME_PAD = 80.0                      # diagram px of slack around the content bounds


def _rgb(value: str | None) -> RGBColor | None:
    v = normalize_color(value) if value else None
    return RGBColor.from_string(v[1:]) if v else None


# --------------------------------------------------------------------------- elements
@dataclass
class Element:
    """One PowerPoint object, positioned in frame coordinates (diagram px)."""
    kind: str                          # rect | picture | line | text
    x: float
    y: float
    w: float
    h: float
    name: str = ""
    fill: str | None = None
    stroke: str | None = None
    stroke_w: float = 1.0
    dashed: bool = False
    radius: float = 0.0                # corner radius in diagram px
    image: Path | None = None
    points: list[tuple[float, float]] = field(default_factory=list)
    head_arrow: bool = False
    tail_arrow: bool = False
    arrowheads: list[list[tuple[float, float]]] = field(default_factory=list)
    label: Label | None = None


@dataclass
class Result:
    path: Path
    slide_index: int
    counts: dict
    frame: tuple[float, float, float, float]
    content: tuple[float, float, float, float]
    workdir: Path | None = None


# --------------------------------------------------------------------------- helpers
def _label_name(cell: model.Cell, kind: str) -> str:
    text = strip_markup(cell.value).strip().replace("\n", " ")[:40]
    return f"{kind}: {text}" if text else f"{kind} {cell.id[-4:]}"


_DATA_URI = re.compile(r"data:image/(png|jpeg|jpg|gif),([A-Za-z0-9+/=]+)$")


def extract_bitmap(cell: model.Cell, workdir: Path) -> Path | None:
    """Decode `image=data:...`, honouring a separate `clipPath=inset(t r b l)` crop.

    The crop lives in its own style key, not inside the data URI — miss it and
    cropped logos come out shrunken.
    """
    img = cell.st.get("image")
    if not isinstance(img, str):
        return None
    m = _DATA_URI.match(img)
    if not m:
        return None
    ext = "png" if m.group(1) == "png" else "jpg"
    dest = workdir / f"img_{cell.id}.{ext}"
    if dest.exists():
        return dest
    raw = base64.b64decode(m.group(2))
    clip = re.match(r"inset\(([^)]*)\)", str(cell.st.get("clipPath") or ""))
    if clip:
        from PIL import Image

        vals = [float(v.rstrip("%")) / 100 for v in clip.group(1).split()]
        top, right, bottom, left = (vals + vals[:1] * 4)[:4]
        with Image.open(io.BytesIO(raw)) as im:
            W, H = im.size
            im.crop((round(left * W), round(top * H),
                     round((1 - right) * W), round((1 - bottom) * H))).save(dest)
    else:
        dest.write_bytes(raw)
    return dest


def collect(page: model.Page, svg: SvgMap, rendered: dict[str, stencils.Rendered],
            frame: tuple[float, float, float, float], workdir: Path) -> list[Element]:
    """Walk the page in paint order and describe every object to emit.

    Everything is expressed in SVG coordinates, since that is the space the label
    boxes and edge routes come back in.
    """
    rx, ry, _, _ = svg.frame_rect(model.FRAME_ID)
    fx, fy = frame[0] - rx, frame[1] - ry
    out: list[Element] = []

    for cell in page.cells:
        if cell.id == model.FRAME_ID:
            continue
        label = svg.label(cell.id)

        if cell.is_edge:
            route = svg.edge(cell.id)
            if route:
                out.append(Element(
                    kind="line", x=0, y=0, w=0, h=0, name=_label_name(cell, "line"),
                    stroke=route.color, stroke_w=route.width, points=route.points,
                    arrowheads=route.arrowheads,
                    head_arrow=cell.st.get("startArrow") not in (None, "none"),
                    tail_arrow=cell.st.get("endArrow") not in (None, "none"),
                ))
            if label:
                out.append(Element(kind="text", x=label.x, y=label.y, w=label.w, h=label.h,
                                   name=_label_name(cell, "label"), label=label))
            continue

        if cell.box is None:
            continue
        bx, by = cell.x - fx, cell.y - fy

        if cell.is_group:
            # native border + separately rendered corner badge, so both stay editable
            out.append(Element(kind="rect", x=bx, y=by, w=cell.w, h=cell.h,
                               name=_label_name(cell, "area"),
                               stroke=str(cell.st.get("strokeColor") or "#232F3E"), stroke_w=1.0,
                               dashed=cell.st.get("dashed") == "1"))
            r = rendered.get(cell.id)
            if r:
                out.append(Element(kind="picture", x=r.x - fx, y=r.y - fy, w=r.w, h=r.h,
                                   name=_label_name(cell, "area icon"), image=r.path))
        elif cell.id in rendered:
            r = rendered[cell.id]
            out.append(Element(kind="picture", x=r.x - fx, y=r.y - fy, w=r.w, h=r.h,
                               name=_label_name(cell, "icon"), image=r.path))
        elif cell.shape == "image":
            path = extract_bitmap(cell, workdir)
            if path:
                out.append(Element(kind="picture", x=bx, y=by, w=cell.w, h=cell.h,
                                   name=_label_name(cell, "icon"), image=path))
        elif cell.is_text_only:
            pass                                   # the label below is the whole object
        else:
            # draw.io defaults an absent fillColor to white and strokeColor to black,
            # so ask the SVG what it painted instead of reading the style dict
            paint = svg.paint(cell.id)
            out.append(Element(
                kind="rect", x=bx, y=by, w=cell.w, h=cell.h,
                name=_label_name(cell, "area"),
                fill=paint.fill if paint else str(cell.st.get("fillColor") or "") or None,
                stroke=paint.stroke if paint else str(cell.st.get("strokeColor") or "") or None,
                stroke_w=(paint.stroke_width if paint
                          else float(cell.st.get("strokeWidth", 1) or 1)),
                dashed=paint.dashed if paint else cell.st.get("dashed") == "1",
                radius=paint.radius if paint else 0.0,
            ))

        if label:
            out.append(Element(kind="text", x=label.x, y=label.y, w=label.w, h=label.h,
                               name=_label_name(cell, "label"), label=label))
    return out


def content_bounds(elements: list[Element]) -> tuple[float, float, float, float]:
    xs, ys = [], []
    for e in elements:
        if e.kind == "line":
            pts = [*e.points, *(p for poly in e.arrowheads for p in poly)]
            xs += [p[0] for p in pts]
            ys += [p[1] for p in pts]
        else:
            xs += [e.x, e.x + e.w]
            ys += [e.y, e.y + e.h]
    if not xs:
        raise ValueError("nothing to place on the slide")
    return min(xs), min(ys), max(xs) - min(xs), max(ys) - min(ys)


# --------------------------------------------------------------------------- emitting
class _Emitter:
    def __init__(self, shapes, scale: float, off_x: int, off_y: int, origin, font, ea_font):
        self.shapes = shapes
        self.k = scale                       # EMU per diagram px
        self.off_x, self.off_y = off_x, off_y
        self.ox, self.oy = origin
        self.font = font
        self.ea_font = ea_font

    def X(self, v):
        return Emu(int(round(self.off_x + (v - self.ox) * self.k)))

    def Y(self, v):
        return Emu(int(round(self.off_y + (v - self.oy) * self.k)))

    def L(self, v):
        return Emu(max(1, int(round(v * self.k))))

    def pt(self, px):
        return round(px * self.k / EMU_PER_PT, 1)

    @staticmethod
    def _no_theme_effects(shape):
        # p:style carries effectRef -> a theme drop shadow on every autoshape/freeform
        shape.shadow.inherit = False

    def _dash(self, line):
        ln = line._get_or_add_ln()
        for old in ln.findall(qn("a:prstDash")):
            ln.remove(old)
        ln.append(ln.makeelement(qn("a:prstDash"), {"val": "dash"}))

    def _arrows(self, line, head: bool, tail: bool):
        ln = line._get_or_add_ln()
        for tag, on in ((qn("a:headEnd"), head), (qn("a:tailEnd"), tail)):
            if on:
                ln.append(ln.makeelement(tag, {"type": "triangle", "w": "med", "len": "med"}))

    def rect(self, e: Element):
        shape = MSO_SHAPE.ROUNDED_RECTANGLE if e.radius else MSO_SHAPE.RECTANGLE
        sh = self.shapes.add_shape(shape, self.X(e.x), self.Y(e.y),
                                   self.L(e.w), self.L(e.h))
        sh.name = e.name
        if e.radius and min(e.w, e.h):
            # roundRect's adjustment is the radius as a fraction of the shorter side
            sh.adjustments[0] = min(0.5, e.radius / min(e.w, e.h))
        fill = _rgb(e.fill) if e.fill and e.fill != "none" else None
        if fill:
            sh.fill.solid()
            sh.fill.fore_color.rgb = fill
        else:
            sh.fill.background()
        stroke = _rgb(e.stroke) if e.stroke and e.stroke != "none" else None
        if stroke:
            sh.line.color.rgb = stroke
            sh.line.width = self.L(e.stroke_w)
            if e.dashed:
                self._dash(sh.line)
        else:
            sh.line.fill.background()
        self._no_theme_effects(sh)
        sh.text_frame.text = ""
        return sh

    def picture(self, e: Element):
        pic = self.shapes.add_picture(str(e.image), self.X(e.x), self.Y(e.y),
                                      self.L(e.w), self.L(e.h))
        pic.name = e.name
        return pic

    def line(self, e: Element):
        """One freeform holding the route plus draw.io's own arrowhead polygons.

        Keeping the arrowheads inside the same shape means the connector still moves as
        one object, while the tips stay the size draw.io drew them at instead of shrinking
        with the line width the way PowerPoint's own arrowheads do.
        """
        pts = e.points
        ff = self.shapes.build_freeform(self.X(pts[0][0]), self.Y(pts[0][1]))
        ff.add_line_segments([(self.X(x), self.Y(y)) for x, y in pts[1:]], close=False)
        sh = ff.convert_to_shape()
        sh.name = e.name
        colour = _rgb(e.stroke) or RGBColor(0, 0, 0)
        sh.line.color.rgb = colour
        sh.line.width = self.L(e.stroke_w)
        if e.arrowheads:
            sh.fill.solid()                       # only the arrowhead subpaths take it
            sh.fill.fore_color.rgb = colour
            self._rebuild_with_arrowheads(sh, pts, e.arrowheads)
        else:
            sh.fill.background()
            self._arrows(sh.line, e.head_arrow, e.tail_arrow)
        self._no_theme_effects(sh)
        return sh

    def _rebuild_with_arrowheads(self, shape, route, polygons):
        """Redo the custGeom over route + arrowheads, so no tip falls outside the bounds."""
        emu = [[(int(self.X(x)), int(self.Y(y))) for x, y in poly]
               for poly in [route, *polygons]]
        xs = [p[0] for poly in emu for p in poly]
        ys = [p[1] for poly in emu for p in poly]
        left, top = min(xs), min(ys)
        w, h = max(xs) - left, max(ys) - top

        shape.left, shape.top, shape.width, shape.height = Emu(left), Emu(top), Emu(w), Emu(h)
        path_lst = shape._element.spPr.find(qn("a:custGeom")).find(qn("a:pathLst"))
        for old in list(path_lst):
            path_lst.remove(old)

        for i, poly in enumerate(emu):
            closed = i > 0                                   # index 0 is the route itself
            path = path_lst.makeelement(
                qn("a:path"), {"w": str(w), "h": str(h), "fill": "norm" if closed else "none"})
            for j, (px, py) in enumerate(poly):
                node = path.makeelement(qn("a:moveTo") if j == 0 else qn("a:lnTo"), {})
                pt = node.makeelement(qn("a:pt"), {"x": str(px - left), "y": str(py - top)})
                node.append(pt)
                path.append(node)
            if closed:
                path.append(path.makeelement(qn("a:close"), {}))
            path_lst.append(path)

    def text(self, e: Element):
        lab = e.label
        tb = self.shapes.add_textbox(self.X(e.x), self.Y(e.y), self.L(e.w), self.L(e.h))
        tb.name = e.name
        tf = tb.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bg = _rgb(lab.bg)
        if bg:
            tb.fill.solid()
            tb.fill.fore_color.rgb = bg
        else:
            tb.fill.background()
        tb.line.fill.background()
        for i, line in enumerate(lab.lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER if lab.align == "center" else PP_ALIGN.LEFT
            # a line may mix sizes (a 20px emoji next to 12px text); space it by the tallest
            p.line_spacing = Pt(self.pt(max(r.size for r in line) * 1.25))
            for spec in line:
                run = p.add_run()
                run.text = spec.text
                run.font.size = Pt(self.pt(spec.size))
                run.font.bold = spec.bold
                run.font.name = self.font or lab.family
                colour = _rgb(spec.color)
                if colour:
                    run.font.color.rgb = colour
                if self.ea_font:
                    rPr = run._r.get_or_add_rPr()
                    rPr.append(rPr.makeelement(qn("a:ea"), {"typeface": self.ea_font}))
        return tb

    def emit(self, elements: list[Element]) -> dict:
        counts = {"rect": 0, "picture": 0, "line": 0, "text": 0}
        for e in elements:
            getattr(self, e.kind)(e)
            counts[e.kind] += 1
        return counts


# --------------------------------------------------------------------------- slides
def _blank_layout(prs):
    """The layout with the fewest placeholders — 'Blank' in stock templates."""
    layouts = list(prs.slide_layouts)
    return min(layouts, key=lambda lay: len(lay.placeholders)) if layouts else None


def _target_slide(prs, slide_no: int | None, replace: bool):
    if slide_no is None:
        slide = prs.slides.add_slide(_blank_layout(prs))
        return slide, len(prs.slides)
    if not 1 <= slide_no <= len(prs.slides):
        raise ValueError(f"--slide {slide_no} is out of range (deck has {len(prs.slides)} slides)")
    slide = prs.slides[slide_no - 1]
    if replace:
        for sh in list(slide.shapes):
            sh._element.getparent().remove(sh._element)
    return slide, slide_no


# --------------------------------------------------------------------------- entry point
def convert(source: str | Path, output: str | Path | None = None, *,
            into: str | Path | None = None, slide: int | None = None, replace: bool = False,
            page: int = 1, slide_size: str = "16:9", margin: float = 0.0,
            scale: int = 6, drawio: str | None = None, font: str | None = None,
            ea_font: str | None = None, keep_workdir: bool = False,
            progress=lambda *_: None) -> Result:
    """Convert one page of `source` into native shapes on a PowerPoint slide."""
    source = Path(source).expanduser().resolve()
    if not source.exists():
        raise FileNotFoundError(f"no such diagram: {source}")
    binary = stencils.find_drawio(drawio)

    pages = model.parse(source)
    if not 1 <= page <= len(pages):
        names = ", ".join(f"{p.index}:{p.name}" for p in pages)
        raise ValueError(f"--page {page} is out of range. Pages: {names}")
    pg = pages[page - 1]

    workdir = Path(tempfile.mkdtemp(prefix="drawio2pptx-"))
    try:
        src_text = source.read_text(encoding="utf-8")
        bx, by, bw, bh = pg.content_bbox()
        frame = (bx - FRAME_PAD, by - FRAME_PAD, bw + 2 * FRAME_PAD, bh + 2 * FRAME_PAD)
        framed = workdir / "framed.drawio"
        framed.write_text(model.with_frame(src_text, page, frame), encoding="utf-8")
        progress("exporting reference SVG")
        svg = SvgMap.load(stencils.export_svg(binary, framed, workdir / "page.svg", page))

        rendered = stencils.render_cells(binary, pg, frame, workdir, scale=scale,
                                         progress=progress)
        elements = collect(pg, svg, rendered, frame, workdir)
        cx, cy, cw, ch = content_bounds(elements)

        progress("writing slide")
        if into:
            prs = Presentation(str(Path(into).expanduser()))
            out_path = Path(output).expanduser() if output else Path(into).expanduser()
        else:
            prs = Presentation()
            if slide_size in SLIDE_SIZES:
                prs.slide_width, prs.slide_height = SLIDE_SIZES[slide_size]
            elif slide_size == "auto":
                prs.slide_width = 12192000
                prs.slide_height = int(12192000 * ch / cw)
            else:
                w, _, h = slide_size.partition("x")
                prs.slide_width, prs.slide_height = int(float(w) * 914400), int(float(h) * 914400)
            slide = None
            out_path = Path(output).expanduser() if output else source.with_suffix(".pptx")

        target, index = _target_slide(prs, slide, replace)
        usable_w = prs.slide_width * (1 - 2 * margin)
        usable_h = prs.slide_height * (1 - 2 * margin)
        k = min(usable_w / cw, usable_h / ch)
        off_x = int((prs.slide_width - cw * k) / 2)
        off_y = int((prs.slide_height - ch * k) / 2)

        counts = _Emitter(target.shapes, k, off_x, off_y, (cx, cy), font, ea_font).emit(elements)
        out_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(out_path))
        return Result(path=out_path, slide_index=index, counts=counts, frame=frame,
                      content=(cx, cy, cw, ch),
                      workdir=workdir if keep_workdir else None)
    finally:
        if not keep_workdir:
            shutil.rmtree(workdir, ignore_errors=True)
