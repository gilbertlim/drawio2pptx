"""Eyeball check: redraw the generated slide and stack it against draw.io's own export.

This is deliberately an *approximate* renderer — it exists to catch shapes that landed
in the wrong place, not to be a PowerPoint clone. Text metrics will differ slightly
because PIL is not laying out text the way PowerPoint does.
"""
from __future__ import annotations

import io
import math
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

_FONT_CANDIDATES = [
    "/System/Library/Fonts/AppleSDGothicNeo.ttc",
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    "C:/Windows/Fonts/malgun.ttf",
    "C:/Windows/Fonts/arial.ttf",
]


_BOLD_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",   # no bold cut; handled below
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "C:/Windows/Fonts/malgunbd.ttf",
    "C:/Windows/Fonts/arialbd.ttf",
]


def _load_font(size: int, bold: bool = False):
    """Best available face at `size`. Bold is resolved by scanning .ttc collections."""
    from PIL import ImageFont

    if bold:
        for path in _FONT_CANDIDATES:
            if not path.endswith(".ttc") or not Path(path).exists():
                continue
            for index in range(10):
                try:
                    font = ImageFont.truetype(path, size, index=index)
                except Exception:
                    break
                if "bold" in (font.font.style or "").lower():
                    return font
        for path in _BOLD_CANDIDATES:
            if Path(path).exists():
                try:
                    return ImageFont.truetype(path, size)
                except Exception:
                    continue
    for path in _FONT_CANDIDATES:
        if Path(path).exists():
            try:
                return ImageFont.truetype(path, size)
            except Exception:
                continue
    return ImageFont.load_default()


def _srgb(el, path) -> tuple | None:
    node = el.find(path) if el is not None else None
    if node is None:
        return None
    clr = node.find(qn("a:srgbClr"))
    if clr is None:
        return None
    v = clr.get("val")
    return tuple(int(v[i:i + 2], 16) for i in (0, 2, 4))


def render_slide(pptx_path: str | Path, slide_index: int, out_png: str | Path,
                 width_px: int = 2000) -> Path:
    """Draw slide `slide_index` (1-based) of a deck to a PNG."""
    from PIL import Image, ImageDraw

    prs = Presentation(str(pptx_path))
    slide = prs.slides[slide_index - 1]
    scale = width_px / prs.slide_width                     # px per EMU
    W = width_px
    H = max(1, round(prs.slide_height * scale))
    img = Image.new("RGB", (W, H), "white")
    dr = ImageDraw.Draw(img, "RGBA")
    fonts: dict[tuple, object] = {}

    for shape in slide.shapes:
        el = shape._element
        spPr = el.find(qn("p:spPr"))
        x0, y0 = shape.left * scale, shape.top * scale
        x1, y1 = (shape.left + shape.width) * scale, (shape.top + shape.height) * scale

        if shape.shape_type == 13:                          # picture
            blip = el.find(qn("p:blipFill")).find(qn("a:blip"))
            part = shape.part.related_part(blip.get(qn("r:embed")))
            with Image.open(io.BytesIO(part.blob)) as src:
                pic = src.convert("RGBA").resize(
                    (max(1, round(x1 - x0)), max(1, round(y1 - y0))), Image.LANCZOS)
            img.paste(pic, (round(x0), round(y0)), pic)
            continue

        fill = _srgb(spPr, qn("a:solidFill"))
        ln = spPr.find(qn("a:ln")) if spPr is not None else None
        stroke = _srgb(ln, qn("a:solidFill")) if ln is not None else None
        lw = max(1, round(int(ln.get("w", 9525)) * scale)) if ln is not None else 1
        dashed = ln is not None and ln.find(qn("a:prstDash")) is not None

        geom = spPr.find(qn("a:custGeom")) if spPr is not None else None
        if geom is not None:
            # a connector is one open subpath for the route plus closed, filled subpaths
            # for the arrowheads draw.io drew
            for path in geom.find(qn("a:pathLst")).findall(qn("a:path")):
                pw, ph = int(path.get("w") or 0), int(path.get("h") or 0)
                pts = []
                for node in path:
                    pt = node.find(qn("a:pt"))
                    if pt is None:
                        continue
                    px, py = int(pt.get("x")), int(pt.get("y"))
                    pts.append((x0 + (px / pw * (x1 - x0) if pw else 0),
                                y0 + (py / ph * (y1 - y0) if ph else 0)))
                if len(pts) < 2:
                    continue
                if path.get("fill") == "norm":
                    dr.polygon(pts, fill=fill or stroke or (0, 0, 0))
                    continue
                dr.line(pts, fill=stroke or (0, 0, 0), width=lw)
                for tag, tip, prev in ((qn("a:headEnd"), pts[0], pts[1]),
                                       (qn("a:tailEnd"), pts[-1], pts[-2])):
                    if ln is not None and ln.find(tag) is not None:
                        ang = math.atan2(tip[1] - prev[1], tip[0] - prev[0])
                        size = max(4.0, lw * 4)
                        dr.polygon([tip,
                                    (tip[0] - size * math.cos(ang - 0.4),
                                     tip[1] - size * math.sin(ang - 0.4)),
                                    (tip[0] - size * math.cos(ang + 0.4),
                                     tip[1] - size * math.sin(ang + 0.4))],
                                   fill=stroke or (0, 0, 0))
            continue

        if fill:
            dr.rectangle([x0, y0, x1, y1], fill=fill)
        if stroke:
            if dashed:
                _dashed_rect(dr, x0, y0, x1, y1, stroke, lw, max(4.0, 12 * scale * 9525))
            else:
                dr.rectangle([x0, y0, x1, y1], outline=stroke, width=lw)

        if shape.has_text_frame and shape.text_frame.text.strip():
            _draw_text(dr, shape, x0, y0, x1, y1, scale, fonts)

    out_png = Path(out_png)
    img.save(out_png)
    return out_png


def _dashed_rect(dr, x0, y0, x1, y1, color, lw, seg):
    for ax, ay, bx, by in ((x0, y0, x1, y0), (x1, y0, x1, y1),
                           (x1, y1, x0, y1), (x0, y1, x0, y0)):
        length = math.hypot(bx - ax, by - ay)
        steps = max(1, int(length / max(seg, 1)))
        for i in range(0, steps, 2):
            t0, t1 = i / steps, min(1.0, (i + 1) / steps)
            dr.line([ax + (bx - ax) * t0, ay + (by - ay) * t0,
                     ax + (bx - ax) * t1, ay + (by - ay) * t1], fill=color, width=lw)


def _draw_text(dr, shape, x0, y0, x1, y1, scale, fonts):
    lines = []
    for para in shape.text_frame.paragraphs:
        if not para.runs:
            continue
        run = para.runs[0]
        size_pt = run.font.size.pt if run.font.size else 12
        px = max(6, round(size_pt * 12700 * scale))
        key = (px, bool(run.font.bold))
        if key not in fonts:
            fonts[key] = _load_font(px, bool(run.font.bold))
        colour = (0, 0, 0)
        if run.font.color and run.font.color.type is not None:
            try:
                hexval = str(run.font.color.rgb)
                colour = tuple(int(hexval[i:i + 2], 16) for i in (0, 2, 4))
            except Exception:
                pass
        lh = (para.line_spacing.pt * 12700 * scale) if para.line_spacing else px * 1.25
        lines.append(("".join(r.text for r in para.runs), fonts[key], colour, lh,
                      str(para.alignment or "")))
    if not lines:
        return
    cy = (y0 + y1) / 2 - sum(line[3] for line in lines) / 2
    for text, font, colour, lh, align in lines:
        bb = dr.textbbox((0, 0), text, font=font)
        tx = x0 if align.startswith("LEFT") else (x0 + x1) / 2 - (bb[2] - bb[0]) / 2
        dr.text((tx, cy + lh / 2), text, font=font, fill=colour, anchor="lm")
        cy += lh


def compare(pptx_path: str | Path, slide_index: int, reference_png: str | Path,
            out_png: str | Path, width_px: int = 2000) -> Path:
    """Stack draw.io's export above our re-render of the slide, plus a difference map."""
    from PIL import Image, ImageChops

    mine = Image.open(render_slide(pptx_path, slide_index,
                                   Path(out_png).with_suffix(".slide.png"), width_px)).convert("RGB")
    with Image.open(reference_png) as r:
        ref = r.convert("RGB")
    ref = _fit(ref, mine.size)
    diff = ImageChops.difference(ref, mine).convert("L").point(lambda v: 255 if v > 60 else 0)
    gap = 10
    canvas = Image.new("RGB", (mine.width, mine.height * 3 + gap * 2), "#c0392b")
    canvas.paste(ref, (0, 0))
    canvas.paste(mine, (0, mine.height + gap))
    canvas.paste(diff.convert("RGB"), (0, (mine.height + gap) * 2))
    out_png = Path(out_png)
    canvas.save(out_png)
    return out_png


def _fit(img, size):
    """Letterbox `img` onto a white canvas of `size`, preserving aspect."""
    from PIL import Image

    W, H = size
    k = min(W / img.width, H / img.height)
    scaled = img.resize((max(1, round(img.width * k)), max(1, round(img.height * k))),
                        Image.LANCZOS)
    canvas = Image.new("RGB", size, "white")
    canvas.paste(scaled, ((W - scaled.width) // 2, (H - scaled.height) // 2))
    return canvas
