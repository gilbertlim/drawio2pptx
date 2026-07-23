"""End-to-end conversion. Skipped when draw.io desktop is not installed."""
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from drawio2pptx import build, convert, stencils

SAMPLE = Path(__file__).parent.parent / "examples" / "sample.drawio"

try:
    stencils.find_drawio()
    HAVE_DRAWIO = True
except stencils.DrawioNotFound:
    HAVE_DRAWIO = False

pytestmark = pytest.mark.skipif(not HAVE_DRAWIO, reason="draw.io desktop is not installed")


@pytest.fixture(scope="module")
def deck(tmp_path_factory):
    out = tmp_path_factory.mktemp("out") / "sample.pptx"
    return convert(SAMPLE, out), out


def test_writes_a_deck_with_one_slide(deck):
    result, out = deck
    assert out.exists()
    prs = Presentation(str(out))
    assert len(prs.slides) == 1
    assert result.slide_index == 1


def test_every_element_kind_is_produced(deck):
    result, _ = deck
    c = result.counts
    assert c["rect"] >= 1, "container rectangles"
    assert c["line"] >= 1, "connectors"
    assert c["text"] >= 1, "labels"
    assert c["picture"] >= 1, "icons"


def test_shapes_are_separate_objects_not_one_image(deck):
    _, out = deck
    shapes = list(Presentation(str(out)).slides[0].shapes)
    assert len(shapes) > 10
    pictures = [s for s in shapes if s.shape_type == 13]
    assert len(pictures) < len(shapes), "must not collapse into a single picture"


def test_labels_keep_their_text_and_weight(deck):
    _, out = deck
    texts = {}
    for shape in Presentation(str(out)).slides[0].shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            run = shape.text_frame.paragraphs[0].runs[0]
            texts[shape.text_frame.text.strip()] = bool(run.font.bold)
    assert "DMZ" in texts and texts["DMZ"] is True, "container titles are bold in the source"
    assert "DBMS" in texts and texts["DBMS"] is False


def test_no_shape_inherits_a_theme_shadow(deck):
    from lxml import etree

    _, out = deck
    for shape in Presentation(str(out)).slides[0].shapes:
        xml = etree.tostring(shape._element)
        if b"effectRef" in xml:
            assert b"effectLst" in xml, f"{shape.name} would pick up the theme drop shadow"


def test_everything_lands_inside_the_slide(deck):
    _, out = deck
    prs = Presentation(str(out))
    for shape in prs.slides[0].shapes:
        assert shape.left >= -1000 and shape.top >= -1000
        assert shape.left + shape.width <= prs.slide_width + 1000
        assert shape.top + shape.height <= prs.slide_height + 1000


def test_margin_shrinks_the_drawing(tmp_path):
    def extent(margin):
        out = tmp_path / f"m{margin}.pptx"
        convert(SAMPLE, out, margin=margin)
        shapes = list(Presentation(str(out)).slides[0].shapes)
        left = min(s.left for s in shapes)
        right = max(s.left + s.width for s in shapes)
        return right - left

    assert extent(0.2) < extent(0.0)


def test_insert_into_an_existing_deck(tmp_path):
    base = tmp_path / "base.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.slides.add_slide(prs.slide_layouts[6])
    prs.save(str(base))

    result = convert(SAMPLE, into=base, slide=2, replace=True)
    out = Presentation(str(base))
    assert len(out.slides) == 2, "must not add a slide when targeting one"
    assert result.slide_index == 2
    assert len(out.slides[1].shapes) == sum(result.counts.values())


def test_append_a_slide_when_no_slide_number_given(tmp_path):
    base = tmp_path / "base.pptx"
    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[5])
    prs.save(str(base))

    result = convert(SAMPLE, into=base)
    assert len(Presentation(str(base)).slides) == 2
    assert result.slide_index == 2


def test_out_of_range_slide_is_rejected(tmp_path):
    base = tmp_path / "base.pptx"
    Presentation().save(str(base))
    with pytest.raises(ValueError, match="out of range"):
        convert(SAMPLE, into=base, slide=9)


def test_out_of_range_page_is_rejected(tmp_path):
    with pytest.raises(ValueError, match="out of range"):
        convert(SAMPLE, tmp_path / "x.pptx", page=99)


def test_slide_size_presets(tmp_path):
    convert(SAMPLE, tmp_path / "a.pptx", slide_size="4:3")
    prs = Presentation(str(tmp_path / "a.pptx"))
    assert (prs.slide_width, prs.slide_height) == build.SLIDE_SIZES["4:3"]

    convert(SAMPLE, tmp_path / "b.pptx", slide_size="13.333x7.5")
    prs = Presentation(str(tmp_path / "b.pptx"))
    assert prs.slide_width == Emu(int(13.333 * 914400))


def test_arrowheads_are_drawn_at_drawio_size(tmp_path):
    """PowerPoint scales its own arrowheads with line width, which erases them on
    hairline connectors. Each arrow must instead carry draw.io's own filled outline."""
    from lxml import etree
    from pptx.oxml.ns import qn

    src = tmp_path / "arrow.drawio"
    src.write_text(
        '<mxfile><diagram name="p" id="p"><mxGraphModel grid="0" page="1">'
        '<root><mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="e" style="endArrow=classic;startArrow=classic;html=1;" edge="1" '
        'parent="1"><mxGeometry relative="1" as="geometry">'
        '<mxPoint x="100" y="100" as="sourcePoint"/>'
        '<mxPoint x="100" y="118" as="targetPoint"/></mxGeometry></mxCell>'
        '<mxCell id="v" value="anchor" style="rounded=0;whiteSpace=wrap;html=1;" vertex="1" '
        'parent="1"><mxGeometry x="200" y="200" width="120" height="60" as="geometry"/>'
        "</mxCell></root></mxGraphModel></diagram></mxfile>",
        encoding="utf-8")

    out = tmp_path / "arrow.pptx"
    convert(src, out)
    connectors = [s for s in Presentation(str(out)).slides[0].shapes
                  if s._element.spPr.find(qn("a:custGeom")) is not None]
    assert connectors, "the edge should have produced a freeform"

    paths = connectors[0]._element.spPr.find(qn("a:custGeom")).find(qn("a:pathLst"))
    kinds = [p.get("fill") for p in paths.findall(qn("a:path"))]
    assert kinds[0] == "none", "the route itself must not be filled"
    assert kinds.count("norm") == 2, f"expected two filled arrowheads, got {kinds}"

    # every arrowhead vertex has to sit inside the shape box or PowerPoint clips the tip
    for path in paths.findall(qn("a:path")):
        w, h = int(path.get("w")), int(path.get("h"))
        for pt in path.iter(qn("a:pt")):
            assert 0 <= int(pt.get("x")) <= w and 0 <= int(pt.get("y")) <= h, \
                etree.tostring(path)


def test_default_styled_shape_keeps_its_fill_and_border(tmp_path):
    """draw.io defaults an absent fillColor to white and strokeColor to black. Reading
    'missing' as 'none' made such shapes vanish, leaving only their label behind."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn

    src = tmp_path / "plain.drawio"
    src.write_text(
        '<mxfile><diagram name="p" id="p"><mxGraphModel grid="0" page="1">'
        '<root><mxCell id="0"/><mxCell id="1" parent="0"/>'
        '<mxCell id="a" value="GW" style="rounded=1;whiteSpace=wrap;html=1;" vertex="1" '
        'parent="1"><mxGeometry x="100" y="100" width="50" height="22" as="geometry"/></mxCell>'
        '<mxCell id="b" value="sq" style="rounded=0;whiteSpace=wrap;html=1;" vertex="1" '
        'parent="1"><mxGeometry x="300" y="100" width="80" height="40" as="geometry"/></mxCell>'
        "</root></mxGraphModel></diagram></mxfile>", encoding="utf-8")

    out = tmp_path / "plain.pptx"
    convert(src, out)
    boxes = {s.name.removeprefix("area: "): s
             for s in Presentation(str(out)).slides[0].shapes if s.name.startswith("area: ")}
    assert sorted(boxes) == ["GW", "sq"]

    for shape in boxes.values():
        spPr = shape._element.spPr
        assert spPr.find(qn("a:noFill")) is None, f"{shape.name} lost its white fill"
        fill = spPr.find(qn("a:solidFill"))
        assert fill is not None and fill.find(qn("a:srgbClr")).get("val") == "FFFFFF"
        ln = spPr.find(qn("a:ln"))
        assert ln is not None and ln.find(qn("a:solidFill")) is not None, \
            f"{shape.name} lost its black border"
        assert ln.find(qn("a:solidFill")).find(qn("a:srgbClr")).get("val") == "000000"

    assert boxes["GW"].auto_shape_type == MSO_SHAPE.ROUNDED_RECTANGLE, \
        "rounded=1 must round the corners"
    assert boxes["sq"].auto_shape_type == MSO_SHAPE.RECTANGLE
