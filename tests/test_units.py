"""Unit tests for the parts that need no draw.io install."""
import base64
import io
from pathlib import Path

import pytest

from drawio2pptx import build, model, stencils
from drawio2pptx.svgmap import SvgMap, normalize_color, strip_markup

SAMPLE = Path(__file__).parent.parent / "examples" / "sample.drawio"


# ------------------------------------------------------------------ style parsing
def test_parse_style_handles_flags_and_values():
    st = model.parse_style("shape=image;html=1;dashed;fillColor=none;")
    assert st["shape"] == "image"
    assert st["dashed"] is True
    assert st["fillColor"] == "none"


def test_parse_style_keeps_parentheses_in_values():
    st = model.parse_style("clipPath=inset(20.41% 20% 20.41% 22.67%);html=1")
    assert st["clipPath"] == "inset(20.41% 20% 20.41% 22.67%)"


# ------------------------------------------------------------------ colours / text
@pytest.mark.parametrize("raw,expected", [
    ("#dae8fc", "#DAE8FC"),
    ("rgb(0, 0, 0)", "#000000"),
    ("light-dark(#999999, #6a6a6a)", "#999999"),
    ("light-dark(rgb(0, 0, 0), rgb(237, 237, 237))", "#000000"),
    ("none", None),
    (None, None),
])
def test_normalize_color(raw, expected):
    assert normalize_color(raw) == expected


def test_strip_markup_turns_divs_into_lines():
    assert strip_markup("LLM<div>call</div>") == "LLM\ncall"
    assert strip_markup("a<br/>b") == "a\nb"
    assert strip_markup("&amp;nbsp;") == "&nbsp;"


# ------------------------------------------------------------------ model
def test_parse_sample_page():
    pages = model.parse(SAMPLE)
    assert len(pages) == 1
    page = pages[0]
    assert page.index == 1
    assert page.cells
    x, y, w, h = page.content_bbox()
    assert w > 0 and h > 0


def test_group_and_text_classification():
    page = model.parse(SAMPLE)[0]
    assert any(c.is_group for c in page.cells), "sample should contain an aws4.group"
    assert all(not (c.is_group and c.is_text_only) for c in page.cells)


def test_with_frame_injects_exactly_one_cell():
    text = SAMPLE.read_text(encoding="utf-8")
    out = model.with_frame(text, 1, (0, 0, 100, 100))
    assert out.count(model.FRAME_ID) == 1
    assert 'width="100" height="100"' in out


def test_with_frame_rejects_missing_page():
    with pytest.raises(ValueError):
        model.with_frame(SAMPLE.read_text(encoding="utf-8"), 5, (0, 0, 10, 10))


# ------------------------------------------------------------------ layer packing
def _cell(cid, x, y, w, h, style="shape=mxgraph.aws4.resourceIcon"):
    return model.Cell(id=cid, value="", style=style, st=model.parse_style(style),
                      is_edge=False, x=x, y=y, w=w, h=h)


def test_disjoint_shapes_share_one_render_pass():
    cells = [_cell("a", 0, 0, 40, 40), _cell("b", 200, 0, 40, 40), _cell("c", 0, 200, 40, 40)]
    assert len(stencils.plan_layers(cells)) == 1


def test_overlapping_shapes_are_split():
    cells = [_cell("a", 0, 0, 40, 40), _cell("b", 10, 10, 40, 40)]
    assert len(stencils.plan_layers(cells)) == 2


def test_group_crop_box_is_just_the_corner_badge():
    group = _cell("g", 0, 0, 400, 300, style="shape=mxgraph.aws4.group;strokeColor=#232F3E")
    x, y, w, h = stencils.crop_box(group)
    assert w <= stencils.GROUP_ICON_EXTENT + 2 * stencils.BLEED
    # a shape sitting inside the group must therefore not collide with it
    inner = _cell("i", 150, 120, 60, 60)
    assert len(stencils.plan_layers([group, inner])) == 1


def test_needs_render_skips_bitmaps_but_keeps_svg_data_uris():
    png = _cell("p", 0, 0, 10, 10, style="shape=image;image=data:image/png,AAAA")
    svg = _cell("s", 0, 0, 10, 10, style="shape=image;image=data:image/svg+xml,PHN2")
    stencil = _cell("t", 0, 0, 10, 10, style="shape=mxgraph.cisco19.rect")
    assert not stencils.needs_render(png)
    assert stencils.needs_render(svg)
    assert stencils.needs_render(stencil)


# ------------------------------------------------------------------ svg reading
TEXT_SVG = """<svg width="100px" height="50px" viewBox="0 0 100 50">
<g data-cell-id="n1"><g><image x="60" y="230" width="60" height="60" href="data:image/png,x"/>
</g><g><g fill="#000000" font-family="Helvetica" text-anchor="middle"
 font-size="12px"><rect fill="#eef0f3" stroke="none" x="60" y="290" width="37" height="15"
 stroke-width="0"/><text x="78" y="300">DBMS</text></g></g></g>
<g data-cell-id="e1"><g transform="translate(0.5,0.5)"><path d="M 10 20 L 40 20 L 40 60"
 fill="none" stroke="#000000" stroke-miterlimit="10"/></g></g>
<g data-cell-id="fr"><g transform="translate(0.5,0.5)"><rect x="3" y="4" width="100"
 height="50" fill="none" stroke="none"/></g></g></svg>"""


def test_label_uses_the_background_rect_as_the_exact_box():
    lab = SvgMap(TEXT_SVG).label("n1")
    assert lab.text == "DBMS"
    assert [r.size for line in lab.lines for r in line] == [12.0]
    assert (lab.x, lab.y, lab.w, lab.h) == (60, 290, 37, 15)
    assert lab.size == 12 and lab.align == "center" and lab.bg == "#EEF0F3"


FO_MIXED = """<svg><g data-cell-id="mix">
<g transform="translate(0.5,0.5)"><rect x="0" y="0" width="60" height="30" fill="#ffffff"/></g>
<g><g><switch><foreignObject>
<div xmlns="http://www.w3.org/1999/xhtml" style="display: flex;">
<div style="box-sizing: border-box; font-size: 0; text-align: center; color: #000000; ">
<div style="display: inline-block; font-size: 12px; font-family: Helvetica;">
<font style="font-size: 20px;">\u23f0</font> Text</div></div></div>
</foreignObject><image x="1" y="2" width="60" height="30"/></switch></g></g></g></svg>"""


def test_a_line_keeps_a_separate_size_per_run():
    """A 20px emoji next to 12px text must not drag the text up to 20px."""
    lab = SvgMap(FO_MIXED).label("mix")
    assert len(lab.lines) == 1
    sizes = [(r.text, r.size) for r in lab.lines[0]]
    assert sizes == [("\u23f0", 20.0), (" Text", 12.0)], sizes
    assert lab.size == 20.0, "line spacing follows the tallest run"


# some stencils letter themselves: kubernetesLabel=1 paints "pod" into the icon, and that
# <text> sits in the drawing group, indistinguishable from a label once the group is flat
LETTERED_SVG = """<svg width="100px" height="100px" viewBox="0 0 100 100">
<g data-cell-id="plain"><g transform="translate(0.5,0.5)">
<path d="M 24 0 L 49 31 L 15 48 Z" fill="#2875e2" stroke="none"/>
<g fill="#ffffff" font-family="Arial, Helvetica" text-anchor="middle" font-size="9.6px">
<text x="25" y="39.8">pod</text></g></g>
<g><g fill="#000000" font-family="Helvetica" text-anchor="middle" font-size="12px">
<text x="25" y="66">WAS</text></g></g></g>
<g data-cell-id="html"><g transform="translate(0.5,0.5)">
<path d="M 124 0 L 149 31 L 115 48 Z" fill="#2875e2" stroke="none"/>
<g fill="#ffffff" font-family="Arial, Helvetica" text-anchor="middle" font-size="9.6px">
<text x="125" y="39.8">pod</text></g></g>
<g><g><switch><foreignObject><div xmlns="http://www.w3.org/1999/xhtml" style="display: flex;">
<div style="box-sizing: border-box; font-size: 0; text-align: center; color: #000000; ">
<div style="display: inline-block; font-size: 12px; font-family: Helvetica;">WAS</div>
</div></div></foreignObject><image x="111" y="55" width="27" height="17"/></switch></g></g></g>
<g data-cell-id="bare"><g transform="translate(0.5,0.5)">
<path d="M 224 0 L 249 31 L 215 48 Z" fill="#2875e2" stroke="none"/>
<g fill="#ffffff" font-family="Arial, Helvetica" text-anchor="middle" font-size="9.6px">
<text x="225" y="39.8">pod</text></g></g></g></svg>"""


@pytest.mark.parametrize("cell_id", ["plain", "html"])
def test_a_stencil_that_letters_itself_does_not_shadow_the_cell_label(cell_id):
    """The icon's own "pod" lettering must not be read as the cell's label."""
    assert SvgMap(LETTERED_SVG).label(cell_id).text == "WAS"


def test_a_lettered_stencil_with_no_label_reports_none():
    """Otherwise the lettering comes back as a text box stamped over the icon."""
    assert SvgMap(LETTERED_SVG).label("bare") is None


def test_edge_route_keeps_every_bend():
    route = SvgMap(TEXT_SVG).edge("e1")
    assert route.points == [(10, 20), (40, 20), (40, 60)]
    assert route.color == "#000000"


def test_frame_rect_reports_where_the_frame_landed():
    assert SvgMap(TEXT_SVG).frame_rect("fr") == (3.0, 4.0, 100.0, 50.0)


def test_svg_size():
    assert SvgMap(TEXT_SVG).size == (100.0, 50.0)


# ------------------------------------------------------------------ error paths
def test_find_drawio_reports_how_to_install():
    with pytest.raises(stencils.DrawioNotFound) as exc:
        stencils.find_drawio("/definitely/not/here")
    assert "/definitely/not/here" in str(exc.value)


# ------------------------------------------------------------------ cli entry
def test_missing_file_explains_the_docs_placeholder(capsys):
    from drawio2pptx import cli

    assert cli.main(["diagram.drawio"]) == 1
    err = capsys.readouterr().err
    assert "no such diagram" in err
    assert "placeholder" in err, "the docs use diagram.drawio; say so instead of an OSError"


def test_missing_file_without_the_placeholder_name(capsys):
    from drawio2pptx import cli

    assert cli.main(["nope.drawio"]) == 1
    err = capsys.readouterr().err
    assert "no such diagram" in err
    assert "placeholder" not in err


# ------------------------------------------------------------------ shape paint
PAINT_SVG = """<svg width="10px" height="10px">
<g data-cell-id="round"><g transform="translate(0.5,0.5)"><rect x="388" y="401" width="50"
 height="21.5" rx="3.23" ry="3.23" fill="#ffffff" stroke="#000000" pointer-events="all"/></g></g>
<g data-cell-id="tinted"><g transform="translate(0.5,0.5)"><rect x="80" y="110" width="490"
 height="130" fill="#dae8fc" stroke="none" pointer-events="all"/></g></g></svg>"""


def test_paint_reads_the_resolved_fill_and_border():
    p = SvgMap(PAINT_SVG).paint("round")
    assert p.fill == "#FFFFFF", "an absent fillColor resolves to white, not to no fill"
    assert p.stroke == "#000000"
    assert p.radius == 3.23


def test_paint_keeps_an_explicit_none_border():
    p = SvgMap(PAINT_SVG).paint("tinted")
    assert p.fill == "#DAE8FC"
    assert p.stroke is None
    assert p.radius == 0.0


# ------------------------------------------------------------------ nested cells
def _nested_file(tmp_path: Path, body: str) -> Path:
    p = tmp_path / "nested.drawio"
    p.write_text('<mxfile><diagram name="p" id="p"><mxGraphModel><root>'
                 '<mxCell id="0"/><mxCell id="1" parent="0"/>'
                 f"{body}</root></mxGraphModel></diagram></mxfile>", encoding="utf-8")
    return p


GROUP = ('<mxCell id="g" value="grp" style="group" vertex="1" parent="1">'
         '<mxGeometry x="400" y="300" width="200" height="120" as="geometry"/></mxCell>')


def test_a_child_of_a_group_is_placed_absolutely(tmp_path):
    """mxGraph stores a child's geometry relative to its parent. Reading it as absolute
    drops the shape somewhere else entirely, with no error to show for it."""
    child = ('<mxCell id="c" style="shape=mxgraph.aws4.resourceIcon" vertex="1" parent="g">'
             '<mxGeometry x="20" y="10" width="40" height="40" as="geometry"/></mxCell>')
    page = model.parse(_nested_file(tmp_path, GROUP + child))[0]
    assert page.by_id("c").box == (420.0, 310.0, 40.0, 40.0)


def test_nesting_two_levels_deep_accumulates_both_offsets(tmp_path):
    inner = ('<mxCell id="i" style="group" vertex="1" parent="g">'
             '<mxGeometry x="10" y="10" width="80" height="80" as="geometry"/></mxCell>'
             '<mxCell id="c" style="shape=mxgraph.aws4.resourceIcon" vertex="1" parent="i">'
             '<mxGeometry x="5" y="5" width="40" height="40" as="geometry"/></mxCell>')
    page = model.parse(_nested_file(tmp_path, GROUP + inner))[0]
    assert page.by_id("c").box == (415.0, 315.0, 40.0, 40.0)


def test_a_group_does_not_stretch_the_content_bounds(tmp_path):
    """A relative child read as absolute inflates the bbox, which moves the frame and
    rescales the whole slide."""
    child = ('<mxCell id="c" style="shape=mxgraph.aws4.resourceIcon" vertex="1" parent="g">'
             '<mxGeometry x="20" y="10" width="40" height="40" as="geometry"/></mxCell>')
    page = model.parse(_nested_file(tmp_path, GROUP + child))[0]
    assert page.content_bbox() == (400.0, 300.0, 200.0, 120.0)


def test_edge_waypoints_inside_a_group_move_with_it(tmp_path):
    edge = ('<mxCell id="e" style="edgeStyle=orthogonalEdgeStyle" edge="1" parent="g">'
            '<mxGeometry relative="1" as="geometry">'
            '<mxPoint x="0" y="0" as="sourcePoint"/>'
            '<mxPoint x="50" y="20" as="targetPoint"/></mxGeometry></mxCell>')
    page = model.parse(_nested_file(tmp_path, GROUP + edge))[0]
    assert page.by_id("e").points == [(400.0, 300.0), (450.0, 320.0)]


def test_an_edge_label_keeps_its_relative_geometry(tmp_path):
    """Edge labels are `relative=1` fractions along the edge, not offsets in px."""
    edge = ('<mxCell id="e" style="edgeStyle=orthogonalEdgeStyle" edge="1" parent="1">'
            '<mxGeometry relative="1" as="geometry"/></mxCell>'
            '<mxCell id="lab" value="x" style="edgeLabel" vertex="1" parent="e">'
            '<mxGeometry x="-0.2" y="1" relative="1" as="geometry"/></mxCell>')
    page = model.parse(_nested_file(tmp_path, edge))[0]
    assert (page.by_id("lab").x, page.by_id("lab").y) == (-0.2, 1.0)


# ------------------------------------------------------------------ embedded bitmaps
def _png_bytes(color="red", size=(100, 100)) -> bytes:
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", size, color).save(buf, "PNG")
    return buf.getvalue()


def _image_cell(uri: str, clip: str | None = None):
    style = f"shape=image;image={uri}" + (f";clipPath={clip}" if clip else "")
    return model.Cell(id="img", value="", style=style, st=model.parse_style(style),
                      is_edge=False, x=0, y=0, w=10, h=10)


def _crop_size(tmp_path, clip: str) -> tuple[int, int]:
    from PIL import Image

    tmp_path.mkdir(parents=True, exist_ok=True)
    uri = "data:image/png," + base64.b64encode(_png_bytes()).decode()
    out = build.extract_bitmap(_image_cell(uri, clip), tmp_path)
    with Image.open(out) as im:
        return im.size


def test_inset_shorthand_follows_the_css_rules(tmp_path):
    """`inset(a b)` is top/bottom then left/right; `inset(a b c)` fills left from right."""
    assert _crop_size(tmp_path / "a", "inset(10% 20% 30% 40%)") == (40, 60)
    assert _crop_size(tmp_path / "b", "inset(10% 20%)") == (60, 80)
    assert _crop_size(tmp_path / "c", "inset(10% 20% 30%)") == (60, 60)
    assert _crop_size(tmp_path / "d", "inset(10%)") == (80, 80)


FRAME_ONLY_SVG = ('<svg width="100px" height="100px"><g data-cell-id="__d2p_frame">'
                  '<rect x="0" y="0" width="100" height="100" fill="none" stroke="none"/>'
                  "</g></svg>")


def test_a_group_wrapper_leaves_no_invisible_rectangle(tmp_path):
    """A plain mxGraph group paints nothing. Emitting it anyway puts an invisible shape on
    the slide that swallows clicks meant for what is inside it."""
    style = "group"
    grp = model.Cell(id="grp", value="", style=style, st=model.parse_style(style),
                     is_edge=False, x=0, y=0, w=60, h=60)
    page = model.Page(id="p", name="p", index=1, cells=[grp])
    els = build.collect(page, SvgMap(FRAME_ONLY_SVG), {}, (0, 0, 100, 100), tmp_path)
    assert [e.name for e in els] == []


def test_an_image_that_cannot_be_decoded_is_reported(tmp_path):
    """An externally linked image has no payload to extract. Emitting nothing for it left
    a hole in the slide with no hint that anything was skipped."""
    cell = _image_cell("https://example.com/logo.png")
    page = model.Page(id="p", name="p", index=1, cells=[cell])
    seen = []
    build.collect(page, SvgMap(FRAME_ONLY_SVG), {}, (0, 0, 100, 100), tmp_path,
                  progress=seen.append)
    assert any(cell.id in m for m in seen), seen


# ------------------------------------------------------------------ curved routes
def test_a_rounded_corner_keeps_its_vertex():
    """draw.io draws rounded bends as `Q`. Skipping the command drops the corner and the
    connector cuts across it diagonally."""
    from drawio2pptx.svgmap import _path_points

    pts = _path_points("M 10 20 L 90 20 Q 100 20 100 30 L 100 90")
    assert pts == [(10.0, 20.0), (90.0, 20.0), (100.0, 30.0), (100.0, 90.0)]


# ------------------------------------------------------------------ early validation
@pytest.fixture
def no_render(monkeypatch):
    """Fail loudly if anything reaches draw.io: these checks must come first."""
    def boom(*a, **k):
        raise AssertionError("rendering started before the arguments were validated")

    monkeypatch.setattr(stencils, "find_drawio", lambda *a, **k: "/fake/drawio")
    monkeypatch.setattr(stencils, "export_svg", boom)
    monkeypatch.setattr(stencils, "render_cells", boom)


def _deck(tmp_path, slides: int = 1) -> Path:
    from pptx import Presentation

    prs = Presentation()
    for _ in range(slides):
        prs.slides.add_slide(prs.slide_layouts[6])
    out = tmp_path / "deck.pptx"
    prs.save(str(out))
    return out


def test_a_missing_target_deck_is_reported_not_traced(tmp_path, no_render):
    with pytest.raises(ValueError, match="no such deck"):
        build.convert(SAMPLE, into=tmp_path / "nope.pptx")


def test_an_out_of_range_slide_is_caught_before_rendering(tmp_path, no_render):
    with pytest.raises(ValueError, match="out of range"):
        build.convert(SAMPLE, into=_deck(tmp_path), slide=9)


def test_a_malformed_slide_size_is_caught_before_rendering(tmp_path, no_render):
    with pytest.raises(ValueError, match="--slide-size"):
        build.convert(SAMPLE, tmp_path / "x.pptx", slide_size="13.333*7.5")


def test_an_impossible_margin_is_rejected(tmp_path, no_render):
    with pytest.raises(ValueError, match="--margin"):
        build.convert(SAMPLE, tmp_path / "x.pptx", margin=0.5)


def test_slide_without_into_is_rejected(capsys):
    """A new deck has no slide to target, so silently dropping the flag hides the mistake."""
    from drawio2pptx import cli

    assert cli.main([str(SAMPLE), "--slide", "3"]) == 2
    assert "--into" in capsys.readouterr().err
